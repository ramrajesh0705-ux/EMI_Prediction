from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ppt_comparison():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "MLflow Runs Comparison Report"
    subtitle.text = "EMI Prediction App - FinTech ML Suite\nGenerated on: April 29, 2026"

    # Slide 2: Classification Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Classification Experiments Overview'
    tf = body_shape.text_frame
    tf.text = 'EMI Eligibility Prediction Models'

    p = tf.add_paragraph()
    p.text = '• 4 different classification models tested'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Focus: Predict loan eligibility (Eligible/High_Risk/Not_Eligible)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Metrics: Accuracy, F1 Score, Precision, Recall, ROC AUC'
    p.level = 1

    # Slide 3: Classification Results Table
    table_slide_layout = prs.slide_layouts[5]  # blank slide
    slide = prs.slides.add_slide(table_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Classification Model Performance'

    # Create table
    rows = 5
    cols = 7
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    for i in range(2, 7):
        table.columns[i].width = Inches(0.8)

    # Headers
    headers = ['Model Type', 'Accuracy', 'F1 Score', 'Precision', 'Recall', 'ROC AUC']
    table.cell(0, 0).text = 'Experiment Name'
    for i, header in enumerate(headers):
        table.cell(0, i+1).text = header

    # Data
    data = [
        ['Logistic Regression v2', '0.7647', '0.6124', '0.6284', '0.7283', '0.9160'],
        ['Logistic Regression v1', '0.8973', '0.5825', '0.5697', '0.5959', '0.9339'],
        ['Random Forest Classifier', '0.8088', '0.6475', '0.6433', '0.7560', 'N/A'],
        ['XGBoost Classifier', '0.8943', '0.7310', '0.7209', '0.7966', 'N/A']
    ]

    experiment_names = [
        'EMI Eligibility Prediction Logistic Regression version 2',
        'EMI Eligibility Prediction Logistic Regression Experiment version 1',
        'EMI Eligibility Prediction Random Forrest Classification experiment',
        'EMI Eligibility Prediction XG Boost Classification experiment'
    ]

    for i, (exp_name, metrics) in enumerate(zip(experiment_names, data)):
        table.cell(i+1, 0).text = exp_name
        for j, metric in enumerate(metrics):
            table.cell(i+1, j+1).text = metric

    # Style the table
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Header row styling
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Slide 4: Classification Summary
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Classification Summary'
    tf = body_shape.text_frame
    tf.text = 'Key Findings'

    p = tf.add_paragraph()
    p.text = 'Best Accuracy: Logistic Regression v1 (0.8973)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Best F1 Score: XGBoost Classifier (0.7310)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Best ROC AUC: Logistic Regression v1 (0.9339)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Logistic Regression shows strong probability calibration'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'XGBoost provides best balance of precision and recall'
    p.level = 1

    # Slide 5: Regression Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Regression Experiments Overview'
    tf = body_shape.text_frame
    tf.text = 'Max Monthly EMI Prediction Models'

    p = tf.add_paragraph()
    p.text = '• 3 different regression models tested'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Focus: Predict maximum safe monthly EMI amount'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Metrics: R² Score, MAE, MAPE, MSE'
    p.level = 1

    # Slide 6: Regression Results Table
    slide = prs.slides.add_slide(table_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Regression Model Performance'

    # Create table
    rows = 4
    cols = 6
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(1.5)
    for i in range(2, 6):
        table.columns[i].width = Inches(0.8)

    # Headers
    headers = ['Model Type', 'R² Score', 'MAE', 'MAPE', 'MSE']
    table.cell(0, 0).text = 'Experiment Name'
    for i, header in enumerate(headers):
        table.cell(0, i+1).text = header

    # Data
    data = [
        ['Random Forest Regressor', '0.9608', '0.1653', '0.0213', '0.0727'],
        ['XGBoost Regressor', '0.9438', '0.2282', '0.0312', '0.1043'],
        ['Linear Regression', '0.7587', '0.5089', '0.0683', '0.4478']
    ]

    experiment_names = [
        'Max Monthly EMI Prediction Random Forest Regressor experiment',
        'Max Monthly EMI Prediction XGBoost Regressor experiment',
        'Max Monthly EMI Prediction Linear Regression experiment'
    ]

    for i, (exp_name, metrics) in enumerate(zip(experiment_names, data)):
        table.cell(i+1, 0).text = exp_name
        for j, metric in enumerate(metrics):
            table.cell(i+1, j+1).text = metric

    # Style the table
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Header row styling
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Slide 7: Regression Summary
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Regression Summary'
    tf = body_shape.text_frame
    tf.text = 'Key Findings'

    p = tf.add_paragraph()
    p.text = 'Best R² Score: Random Forest Regressor (0.9608) - 96.08% variance explained'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Best MAE: Random Forest Regressor (0.1653) - lowest absolute error'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Best MAPE: Random Forest Regressor (0.0213) - 2.13% mean absolute percentage error'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Random Forest significantly outperforms other models'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Linear Regression shows poorest performance - indicates non-linear relationships'
    p.level = 1

    # Slide 8: Key Insights
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Key Insights from All Runs'
    tf = body_shape.text_frame
    tf.text = 'Model Selection & Performance Trends'

    p = tf.add_paragraph()
    p.text = 'Model Selection Recommendations:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Classification: XGBoost Classifier (best F1 score)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Regression: Random Forest Regressor (highest accuracy, lowest error)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Performance Trends:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Tree-based models outperform linear models consistently'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Ensemble methods provide better generalization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Logistic Regression shows good probability calibration (high ROC AUC)'
    p.level = 1

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Conclusion'
    tf = body_shape.text_frame
    tf.text = 'Summary & Next Steps'

    p = tf.add_paragraph()
    p.text = 'This comparison demonstrates a thorough model development process'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Multiple algorithms tested for classification and regression tasks'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'EMI prediction domain shows complex, non-linear relationships'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Tree-based ensemble models are recommended for production deployment'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'MLflow provides excellent experiment tracking and comparison capabilities'
    p.level = 1

    prs.save('mlflow_comparison_report.pptx')
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_ppt_comparison()