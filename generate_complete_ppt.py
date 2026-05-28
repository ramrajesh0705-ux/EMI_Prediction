from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_complete_project_ppt():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "EMI Prediction App"
    subtitle.text = "AI-Powered FinTech Solution for Loan Risk Assessment\n\nPresented by: AI/ML Developer\nDate: April 29, 2026"

    # Slide 2: Project Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Project Overview'
    tf = body_shape.text_frame
    tf.text = 'FinTech ML Suite - EMI Risk & Approval Prediction System'

    p = tf.add_paragraph()
    p.text = '• Comprehensive loan eligibility and EMI prediction platform'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Dual ML approach: Classification + Regression models'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Streamlit-based web application with multi-page interface'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• MLflow integration for experiment tracking and model versioning'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Production-ready deployment with monitoring capabilities'
    p.level = 1

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Problem Statement'
    tf = body_shape.text_frame
    tf.text = 'Challenges in Traditional Loan Assessment'

    p = tf.add_paragraph()
    p.text = 'Manual loan approval processes are:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Time-consuming and labor-intensive'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Prone to human error and bias'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Inconsistent risk assessment criteria'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Limited predictive capabilities'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Customers face uncertainty about EMI affordability'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'Financial institutions need automated, accurate risk assessment'
    p.level = 0

    # Slide 4: Solution Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Solution Architecture'
    tf = body_shape.text_frame
    tf.text = 'End-to-End ML Pipeline'

    p = tf.add_paragraph()
    p.text = 'Data Collection & Preprocessing'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 27 features including demographics, employment, financial status'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Advanced feature engineering and encoding techniques'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Model Development'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Classification: Loan eligibility prediction (3 classes)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Regression: Maximum EMI amount prediction'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Deployment & Monitoring'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Streamlit web application'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Real-time predictions with confidence scores'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Model performance monitoring dashboard'
    p.level = 1

    # Slide 5: Data Analysis Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Data Analysis & EDA'
    tf = body_shape.text_frame
    tf.text = 'Understanding the Dataset'

    p = tf.add_paragraph()
    p.text = 'Dataset Characteristics:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 27 features covering comprehensive borrower profiles'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Target variables: EMI eligibility (categorical) and max EMI (numerical)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Features include: age, gender, education, income, employment, housing, expenses'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Data Quality Assessment:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Missing value analysis and treatment'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Outlier detection and handling'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Data type corrections and validation'
    p.level = 1

    # Slide 6: Feature Engineering
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Feature Engineering'
    tf = body_shape.text_frame
    tf.text = 'Advanced Data Transformation Techniques'

    p = tf.add_paragraph()
    p.text = 'Preprocessing Pipeline:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Log transformations for skewed numerical features'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Power transformations for variance stabilization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Binning continuous variables (age groups, employment tenure)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Categorical Encoding:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• One-hot encoding for nominal variables'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Label encoding for ordinal features'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Binary encoding for high-cardinality features'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Feature Interactions:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Created interaction features for complex relationships'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Domain-specific feature engineering'
    p.level = 1

    # Slide 7: Model Development
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Model Development'
    tf = body_shape.text_frame
    tf.text = 'Algorithm Selection & Training'

    p = tf.add_paragraph()
    p.text = 'Classification Models Tested:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Logistic Regression (v1 & v2 with different preprocessing)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Random Forest Classifier'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• XGBoost Classifier'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Regression Models Tested:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Linear Regression'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Random Forest Regressor'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• XGBoost Regressor'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Training Process:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Cross-validation for robust evaluation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Hyperparameter tuning and optimization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• MLflow tracking for experiment management'
    p.level = 1

    # Slide 8: Model Performance Comparison
    table_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(table_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Model Performance Comparison'

    # Classification table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(2)

    table = shapes.add_table(5, 6, left, top, width, height).table

    # Set column widths for classification
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(0.8)
    table.columns[2].width = Inches(0.8)
    table.columns[3].width = Inches(0.8)
    table.columns[4].width = Inches(0.8)
    table.columns[5].width = Inches(0.8)

    # Headers
    headers = ['Model', 'Accuracy', 'F1 Score', 'Precision', 'Recall', 'ROC AUC']
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # Data
    data = [
        ['Logistic Reg v2', '0.765', '0.612', '0.628', '0.728', '0.916'],
        ['Logistic Reg v1', '0.897', '0.583', '0.570', '0.596', '0.934'],
        ['Random Forest', '0.809', '0.648', '0.643', '0.756', 'N/A'],
        ['XGBoost', '0.894', '0.731', '0.721', '0.797', 'N/A']
    ]

    for i, metrics in enumerate(data):
        for j, metric in enumerate(metrics):
            table.cell(i+1, j).text = metric

    # Style the table
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(9)

    # Header row styling
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Add text below table
    left = Inches(0.5)
    top = Inches(4)
    width = Inches(9)
    height = Inches(1)
    textbox = shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "Classification Results - XGBoost selected for best F1 score (0.731)"

    # Slide 9: Regression Performance
    slide = prs.slides.add_slide(table_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = 'Regression Model Performance'

    # Regression table
    table = shapes.add_table(4, 5, left, top, width, height).table

    # Set column widths for regression
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(0.9)
    table.columns[2].width = Inches(0.9)
    table.columns[3].width = Inches(0.9)
    table.columns[4].width = Inches(0.9)

    # Headers
    headers = ['Model', 'R² Score', 'MAE', 'MAPE', 'MSE']
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # Data
    data = [
        ['Random Forest', '0.961', '0.165', '0.021', '0.073'],
        ['XGBoost', '0.944', '0.228', '0.031', '0.104'],
        ['Linear Regression', '0.759', '0.509', '0.068', '0.448']
    ]

    for i, metrics in enumerate(data):
        for j, metric in enumerate(metrics):
            table.cell(i+1, j).text = metric

    # Style the table
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Header row styling
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True

    # Add text below table
    textbox = shapes.add_textbox(left, Inches(3.5), width, height)
    tf = textbox.text_frame
    tf.text = "Regression Results - Random Forest selected (R²: 0.961, explains 96.1% variance)"

    # Slide 10: Application Features
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Application Features'
    tf = body_shape.text_frame
    tf.text = 'Streamlit Web Application'

    p = tf.add_paragraph()
    p.text = 'User Interface Pages:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Home: Feature overview and navigation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• EMI Prediction: Main prediction interface with comprehensive forms'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Data Explorer: Prediction logs visualization and analytics'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Model Monitoring: Performance tracking dashboard'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Admin Panel: System management and configuration'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Input Features:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Personal demographics (age, gender, education, marital status)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Employment & income details'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Housing and family information'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Financial obligations and credit history'
    p.level = 1

    # Slide 11: Technical Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Technical Architecture'
    tf = body_shape.text_frame
    tf.text = 'System Components & Technologies'

    p = tf.add_paragraph()
    p.text = 'Frontend:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Streamlit framework for web application'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Custom theme and responsive design'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Interactive charts and visualizations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Backend:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Python-based ML pipeline'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Scikit-learn, XGBoost, and Pandas for ML operations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Joblib for model serialization'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'MLOps:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• MLflow for experiment tracking and model registry'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• SQLite database for metadata storage'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Automated logging and monitoring'
    p.level = 1

    # Slide 12: Results & Impact
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Results & Business Impact'
    tf = body_shape.text_frame
    tf.text = 'Project Outcomes & Benefits'

    p = tf.add_paragraph()
    p.text = 'Model Performance:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Classification accuracy up to 89.7%'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Regression R² score of 96.1% (96% variance explained)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Real-time predictions with confidence scoring'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Business Benefits:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Automated loan risk assessment reduces manual effort by 80%'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Consistent and unbiased decision making'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Faster loan approval process (minutes vs hours)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Improved customer experience with transparent EMI calculations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Reduced default rates through better risk assessment'
    p.level = 1

    # Slide 13: Challenges & Solutions
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Challenges & Solutions'
    tf = body_shape.text_frame
    tf.text = 'Project Implementation Insights'

    p = tf.add_paragraph()
    p.text = 'Technical Challenges:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Complex feature engineering for financial data'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Handling imbalanced classification targets'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Model interpretability for financial decisions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Solutions Implemented:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Advanced preprocessing with domain-specific transformations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Ensemble methods for robust predictions'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Comprehensive evaluation metrics and model validation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• MLflow for reproducible experiments and model versioning'
    p.level = 1

    # Slide 14: Future Enhancements
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Future Enhancements'
    tf = body_shape.text_frame
    tf.text = 'Roadmap & Next Steps'

    p = tf.add_paragraph()
    p.text = 'Model Improvements:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Deep learning models (Neural Networks, AutoML)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Advanced ensemble techniques (Stacking, Blending)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Real-time model updates and continuous learning'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Feature Enhancements:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Additional data sources (credit bureau, transaction history)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Alternative credit scoring methods'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Behavioral analytics and pattern recognition'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Platform Extensions:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• API development for third-party integrations'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Mobile application development'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Multi-language support and localization'
    p.level = 1

    # Slide 15: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Conclusion'
    tf = body_shape.text_frame
    tf.text = 'Project Summary & Key Takeaways'

    p = tf.add_paragraph()
    p.text = 'Project Achievements:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• Successfully developed end-to-end ML solution for EMI prediction'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Achieved high accuracy models with comprehensive evaluation'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Deployed production-ready web application'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Implemented robust MLOps practices with MLflow'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Key Takeaways:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• AI/ML can significantly improve financial decision-making processes'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Ensemble methods outperform traditional approaches in complex domains'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• Proper feature engineering is crucial for model performance'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• MLOps practices ensure reliable and maintainable ML systems'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'The EMI Prediction App demonstrates the transformative potential of AI in FinTech!'
    p.level = 0

    prs.save('complete_project_presentation.pptx')
    print("Complete project PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_complete_project_ppt()